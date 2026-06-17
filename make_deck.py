# -*- coding: utf-8 -*-
"""Generate the progress presentation for the LLM-as-MR-BT-Planner paper."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ---------------------------------------------------------------
INK     = RGBColor(0x14, 0x1B, 0x2E)   # near-black navy
NAVY    = RGBColor(0x1F, 0x3A, 0x5F)   # deep blue
ACCENT  = RGBColor(0x2E, 0x86, 0xDE)   # bright blue
ACCENT2 = RGBColor(0x16, 0xA0, 0x85)   # teal/green
AMBER   = RGBColor(0xE6, 0x7E, 0x22)   # amber
MUTED   = RGBColor(0x6B, 0x74, 0x86)   # grey
LIGHT   = RGBColor(0xF2, 0xF5, 0xFA)   # panel bg
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Segoe UI"

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, *rest) in para:
            ital = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.italic = ital
            r.font.name = FONT
    return tb

def header(s, kicker, title, num):
    box(s, 0, 0, SW, Inches(1.18), fill=WHITE)
    box(s, Inches(0.6), Inches(0.34), Inches(0.09), Inches(0.52), fill=ACCENT)
    text(s, Inches(0.85), Inches(0.22), Inches(10.5), Inches(0.9), [
        [(kicker, 12, ACCENT, True)],
        [(title, 25, INK, True)],
    ], space_after=2)
    # page number chip
    text(s, SW - Inches(1.2), Inches(0.40), Inches(0.7), Inches(0.4),
         [[(num, 12, MUTED, True)]], align=PP_ALIGN.RIGHT)
    box(s, Inches(0.85), Inches(1.12), SW - Inches(1.7), Pt(1.4), fill=LIGHT)

def chip(s, x, y, w, label, color):
    c = box(s, x, y, w, Inches(0.34), fill=color, round_=True)
    text(s, x, y, w, Inches(0.34), [[(label, 11, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c

# ===========================================================================
# 1 — TITLE
# ===========================================================================
s = slide(); bg(s, INK)
box(s, 0, 0, SW, SH, fill=INK)
# accent band
box(s, 0, Inches(4.55), SW, Inches(0.06), fill=ACCENT)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.5),
     [[("RESEARCH PROGRESS REVIEW", 14, ACCENT, True)]])
text(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.6), [
    [("LLM as a Multi-Robot", 44, WHITE, True)],
    [("Behavior-Tree Planner", 44, WHITE, True)],
], space_after=4, line_spacing=1.0)
text(s, Inches(0.9), Inches(3.55), Inches(11.4), Inches(1.0),
     [[("Generating, validating, and simulating synchronized "
        "multi-robot behavior trees with a language model as the planner.",
        17, RGBColor(0xC3,0xCD,0xDB), False)]], line_spacing=1.1)
text(s, Inches(0.9), Inches(4.85), Inches(11.4), Inches(1.6), [
    [("Scientific paper  ·  work-in-progress", 14, RGBColor(0x8F,0x9D,0xB3), True)],
    [("8 / 8 planned milestones closed   ·   prototype → full LLM-driven framework",
      14, ACCENT2, True)],
], space_after=6)
text(s, Inches(0.9), SH - Inches(0.85), Inches(11.4), Inches(0.5),
     [[("Status as of 17 June 2026", 12, MUTED, False)]])

# ===========================================================================
# 2 — PROBLEM & MOTIVATION
# ===========================================================================
s = slide(); bg(s)
header(s, "WHY THIS WORK", "The problem we attack", "02")
text(s, Inches(0.85), Inches(1.45), Inches(11.6), Inches(1.2), [
    [("Coordinating ", 17, INK, False),
     ("multiple heterogeneous robots", 17, ACCENT, True),
     (" on one task needs a plan: who does what, in what order, "
      "and where they must wait for each other.", 17, INK, False)],
], line_spacing=1.15)

cards = [
    ("Classic planners", AMBER,
     "Symbolic methods (MRBTP-style back-chaining) synthesize the tree "
     "deterministically — powerful, but hand-built and domain-bound."),
    ("Our question", ACCENT,
     "Can an LLM — given only the instruction and the initial world "
     "state — infer the entire synchronized multi-robot plan itself?"),
    ("The catch", ACCENT2,
     "No hidden checklist, no fixed ordering, no deterministic repair. "
     "The model must discover the causal chains on its own."),
]
cw = Inches(3.85); gap = Inches(0.2); x0 = Inches(0.85); y0 = Inches(2.95)
for i,(t,c,body) in enumerate(cards):
    x = x0 + i*(cw+gap)
    box(s, x, y0, cw, Inches(3.0), fill=LIGHT, round_=True)
    box(s, x, y0, cw, Inches(0.5), fill=c, round_=True)
    box(s, x, y0+Inches(0.25), cw, Inches(0.25), fill=LIGHT)  # square off bottom of header
    text(s, x+Inches(0.25), y0+Inches(0.06), cw-Inches(0.5), Inches(0.45),
         [[(t, 16, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.25), y0+Inches(0.72), cw-Inches(0.5), Inches(2.1),
         [[(body, 14, INK, False)]], line_spacing=1.15)

# ===========================================================================
# 3 — RESEARCH STANCE / CONTRIBUTION
# ===========================================================================
s = slide(); bg(s)
header(s, "CORE IDEA", "Research stance: the LLM is the planner", "03")
box(s, Inches(0.85), Inches(1.5), Inches(11.6), Inches(1.15), fill=NAVY, round_=True)
text(s, Inches(1.15), Inches(1.5), Inches(11.0), Inches(1.15), [
    [("All plan ", 18, WHITE, False), ("structure", 18, WHITE, True, True),
     (" is produced by the language model.  The deterministic code only ",
      18, WHITE, False),
     ("verifies and simulates", 18, ACCENT2, True),
     (" — it never authors or repairs the plan.", 18, WHITE, False)],
], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

rows = [
    ("Pure mode (default)", "Model sees only: prompt + initial state + output schema. "
     "Validator reports WHAT is wrong, never task-specific fixes."),
    ("Assisted mode (ablation)", "Optional dependency hints + candidate-producer "
     "suggestions — a baseline to measure how much scaffolding helps."),
    ("Always-on checks", "Acyclicity, capability match, predicate support, sync "
     "consistency — task-agnostic verification that defines “working.”"),
]
y = Inches(3.0)
for i,(t,b) in enumerate(rows):
    yy = y + i*Inches(1.25)
    box(s, Inches(0.85), yy, Inches(11.6), Inches(1.1), fill=LIGHT, round_=True)
    box(s, Inches(0.85), yy, Inches(0.12), Inches(1.1), fill=ACCENT)
    text(s, Inches(1.2), yy+Inches(0.12), Inches(3.4), Inches(0.9),
         [[(t, 16, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.7), yy+Inches(0.12), Inches(7.5), Inches(0.9),
         [[(b, 14, INK, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

# ===========================================================================
# 4 — METHOD / PIPELINE
# ===========================================================================
s = slide(); bg(s)
header(s, "METHOD", "The pipeline (Algorithm 1)", "04")
steps = [
    ("1", "Load", "Multi-robot scenario:\ninit + goal state,\nobjects, locations,\ncapability libraries"),
    ("2", "Generate", "LLM infers task graph,\nrobot assignments,\nsynchronization,\nper-robot BTs"),
    ("3", "Validate", "Static checks:\nstructure, capabilities,\npredicate support,\nsync consistency"),
    ("4", "Self-correct", "Typed errors + sim\ntrace fed back;\nLLM patches\n(no auto-repair)"),
    ("5", "Simulate", "Tick-based symbolic\nBT execution;\ndeadlock / timeout\ndetection"),
]
n=len(steps); cw=Inches(2.18); gap=Inches(0.14); x0=Inches(0.7); y0=Inches(2.0)
colors=[NAVY, ACCENT, ACCENT2, AMBER, NAVY]
for i,(num,t,b) in enumerate(steps):
    x = x0 + i*(cw+gap)
    box(s, x, y0, cw, Inches(3.1), fill=LIGHT, round_=True)
    box(s, x, y0, cw, Inches(0.62), fill=colors[i], round_=True)
    box(s, x, y0+Inches(0.32), cw, Inches(0.30), fill=LIGHT)
    text(s, x+Inches(0.2), y0+Inches(0.05), cw-Inches(0.4), Inches(0.5),
         [[(num+"   "+t, 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.2), y0+Inches(0.8), cw-Inches(0.4), Inches(2.2),
         [[(b, 12.5, INK, False)]], line_spacing=1.12)
    if i < n-1:
        text(s, x+cw-Inches(0.02), y0+Inches(1.2), Inches(0.18), Inches(0.5),
             [[("›", 22, MUTED, True)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.7), Inches(5.45), Inches(12.0), Inches(1.25), fill=NAVY, round_=True)
text(s, Inches(1.0), Inches(5.45), Inches(11.4), Inches(1.25), [
    [("Loop 2–5 until the plan validates AND simulates to the goal, "
      "or the correction budget is spent → one result file.", 15, WHITE, False)],
    [("Three task-agnostic reliability levers: back-chaining method in the prompt  ·  "
      "best-of-N sampling  ·  two-stage generation (action plan → BTs).",
      13, ACCENT2, True)],
], anchor=MSO_ANCHOR.MIDDLE, space_after=6, line_spacing=1.1)

# ===========================================================================
# 5 — WHAT'S BUILT (system highlights)
# ===========================================================================
s = slide(); bg(s)
header(s, "DELIVERED", "What is built so far", "05")
items = [
    ("Declarative domain", "PDDL-style add/delete capability effects — world "
     "semantics live in the data, not hidden engine conventions."),
    ("Real behavior trees", "Sequence / Fallback / Parallel composites, tick engine "
     "with SUCCESS / FAILURE / RUNNING and reactive memory."),
    ("LLM-as-planner core", "Pure mode by default + 3 reliability levers; OpenAI & "
     "Anthropic providers with automatic fallback."),
    ("Validator & simulator", "20+ typed error classes; deadlock-sound tick simulation "
     "feeding structured feedback to the LLM."),
    ("Visualization", "Self-contained HTML report — Mermaid behavior-tree view + "
     "chronological action-plan timeline."),
    ("Reproducible experiments", "Multi-trial runner: success rate, validity rate, "
     "mean±std correction rounds → CSV / Markdown / JSON."),
    ("Real-robot seam", "BehaviorTree.CPP XML export + ROS execution-backend scaffold "
     "for hardware testing."),
    ("Test suite", "Deterministic, LLM-free pytest over predicates, domain, BTs, "
     "validation, simulation, execution, viz."),
]
cw=Inches(5.75); ch=Inches(1.2); gx=Inches(0.85); gy=Inches(1.55)
hgap=Inches(0.18); vgap=Inches(0.16)
for i,(t,b) in enumerate(items):
    r,c = divmod(i,2)
    x = gx + c*(cw+hgap); y = gy + r*(ch+vgap)
    box(s, x, y, cw, ch, fill=LIGHT, round_=True)
    box(s, x, y, Inches(0.12), ch, fill=ACCENT if c==0 else ACCENT2)
    text(s, x+Inches(0.3), y+Inches(0.12), cw-Inches(0.5), Inches(0.4),
         [[(t, 14.5, NAVY, True)]])
    text(s, x+Inches(0.3), y+Inches(0.52), cw-Inches(0.5), Inches(0.62),
         [[(b, 11.5, INK, False)]], line_spacing=1.05)

# ===========================================================================
# 6 — PROGRESS BY WEEK (from tracker)
# ===========================================================================
s = slide(); bg(s)
header(s, "WHAT WE DID — WEEK BY WEEK", "Progress timeline", "06")
weeks = [
    ("Week 1", "18 – 25 May", ACCENT, "Foundations & scoping", [
        "MRBTP paper analysis",
        "Robot platforms studied (Go2, Z1 arm, Franka)",
        "Algorithm analysis",
        "Research plan & proposal (timeline, contributions)",
    ]),
    ("Week 2", "26 May – 1 Jun", ACCENT2, "Setup & baselines", [
        "Research paper template set up (Overleaf)",
        "Baseline code samples — testing existing algorithms",
    ]),
    ("Week 3", "2 – 8 Jun", AMBER, "Our algorithm", [
        "New algorithm for multi-robot LLM BT generation",
        "LLM-as-planner core: generate → validate → simulate → self-correct",
    ]),
    ("Week 4", "9 – 15 Jun", NAVY, "Hardening & visualization", [
        "Improved generation code (best-of-N, two-stage)",
        "Visual BT simulation + HTML report",
    ]),
]
# vertical timeline rail
railx = Inches(1.25); y0=Inches(1.7); rowh=Inches(1.32)
box(s, railx, y0+Inches(0.15), Pt(2.4), Inches(4.9), fill=LIGHT)
for i,(wk,dates,col,theme,items_) in enumerate(weeks):
    y = y0 + i*rowh
    # node
    box(s, railx-Inches(0.16), y+Inches(0.2), Inches(0.42), Inches(0.42),
        fill=col, round_=True)
    text(s, railx-Inches(0.16), y+Inches(0.2), Inches(0.42), Inches(0.42),
         [[("✔", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # week label
    text(s, Inches(0.55), y+Inches(0.18), Inches(0.7), Inches(0.9),
         [[(wk, 14, col, True)], [(dates, 9.5, MUTED, False)]],
         align=PP_ALIGN.LEFT, space_after=1)
    # card
    cx=Inches(1.95); cw2=Inches(10.4); chh=Inches(1.18)
    box(s, cx, y, cw2, chh, fill=LIGHT, round_=True)
    box(s, cx, y, Inches(0.1), chh, fill=col)
    text(s, cx+Inches(0.3), y+Inches(0.1), cw2-Inches(0.6), Inches(0.4),
         [[(theme, 14.5, NAVY, True)]])
    text(s, cx+Inches(0.3), y+Inches(0.5), cw2-Inches(0.6), Inches(0.65),
         [[("•  "+it, 12, INK, False)] for it in items_],
         line_spacing=1.0, space_after=1)

# ===========================================================================
# 7 — SCENARIOS / EVALUATION
# ===========================================================================
s = slide(); bg(s)
header(s, "EVALUATION SETUP", "Benchmark scenarios & metrics", "07")
# two scenario cards
sc = [
    ("gear_assembly", ACCENT, "3-robot symbolic gear-assembly cell",
     ["go2_z1: opens drawer, stages gear tray & screwdriver, returns tool, closes drawer",
      "franka1: holds and stabilizes the gearbase",
      "franka2: picks/mounts gear, picks screwdriver, fastens the screw"]),
    ("sensor_calibration_cell", ACCENT2, "3-robot, dependency-heavy",
     ["More cross-robot synchronization than gear cell",
      "Calibration, inspection, clamp release",
      "Tool return + drawer closure chains"]),
]
cw=Inches(5.75); x0=Inches(0.85); y0=Inches(1.55)
for i,(name,c,sub,pts) in enumerate(sc):
    x=x0+i*(cw+Inches(0.18))
    box(s, x, y0, cw, Inches(3.1), fill=LIGHT, round_=True)
    box(s, x, y0, cw, Inches(0.7), fill=c, round_=True)
    box(s, x, y0+Inches(0.35), cw, Inches(0.35), fill=LIGHT)
    text(s, x+Inches(0.3), y0+Inches(0.06), cw-Inches(0.6), Inches(0.6),
         [[(name, 17, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.3), y0+Inches(0.82), cw-Inches(0.6), Inches(0.4),
         [[(sub, 13, c, True)]])
    text(s, x+Inches(0.3), y0+Inches(1.25), cw-Inches(0.6), Inches(1.8),
         [[("•  "+p, 12.5, INK, False)] for p in pts], line_spacing=1.1, space_after=6)
# metrics strip
box(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.7), fill=NAVY, round_=True)
text(s, Inches(1.1), Inches(5.1), Inches(11.2), Inches(0.4),
     [[("Reported per trial, aggregated over scenarios × trials:", 13, ACCENT2, True)]])
mets=[("Success rate","goal reached"),("Validity rate","plan passes checks"),
      ("Correction rounds","mean ± std"),("Pure vs assisted","ablation delta")]
mw=Inches(2.75)
for i,(t,b) in enumerate(mets):
    x=Inches(1.1)+i*mw
    text(s, x, Inches(5.55), mw-Inches(0.2), Inches(1.0), [
        [(t, 16, WHITE, True)], [(b, 12, RGBColor(0xC3,0xCD,0xDB), False)],
    ], space_after=2)

# ===========================================================================
# 8 — FURTHER PLANS (paper sections to write)
# ===========================================================================
s = slide(); bg(s)
header(s, "WHAT'S NEXT", "Further plans — writing the paper", "08")
text(s, Inches(0.85), Inches(1.3), Inches(11.6), Inches(0.5),
     [[("The framework is built; the remaining work is the paper itself, "
        "section by section:", 14, MUTED, False)]])
steps=[
    ("Intro (Problem)", AMBER,
     "Motivate synchronized multi-robot planning and the LLM-as-planner question."),
    ("Related works", ACCENT,
     "Position vs. MRBTP-style symbolic planners and LLM-for-robotics literature."),
    ("Background + Methodology", ACCENT2,
     "Domain model, behavior trees, world-state semantics, and evaluation setup."),
    ("Our method / algorithm", NAVY,
     "Self-correction loop, two-stage generation, validator, tick simulator."),
    ("Experiments", AMBER,
     "Run pure vs. assisted across scenarios & trials; ablate the reliability levers."),
    ("Results + Conclusion", ACCENT2,
     "Report success / validity / correction metrics; discuss findings and future work."),
]
cw=Inches(3.82); ch=Inches(1.9); x0=Inches(0.85); y0=Inches(1.95)
hgap=Inches(0.18); vgap=Inches(0.2)
for i,(t,c,b) in enumerate(steps):
    r,col=divmod(i,3)
    x=x0+col*(cw+hgap); y=y0+r*(ch+vgap)
    box(s, x, y, cw, ch, fill=LIGHT, round_=True)
    box(s, x, y, cw, Inches(0.12), fill=c)
    text(s, x+Inches(0.28), y+Inches(0.28), cw-Inches(0.5), Inches(0.5),
         [[(str(i+1), 26, c, True)]])
    text(s, x+Inches(0.28), y+Inches(0.82), cw-Inches(0.5), Inches(0.4),
         [[(t, 15, NAVY, True)]])
    text(s, x+Inches(0.28), y+Inches(1.18), cw-Inches(0.5), Inches(0.6),
         [[(b, 11.5, INK, False)]], line_spacing=1.05)

# ===========================================================================
# 9 — CLOSING
# ===========================================================================
s = slide(); bg(s, INK)
box(s, 0, Inches(2.9), SW, Inches(0.06), fill=ACCENT)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.4), [
    [("Summary", 14, ACCENT, True)],
    [("A working LLM-as-planner framework —", 30, WHITE, True)],
    [("now ready for the experimental study.", 30, WHITE, True)],
], space_after=4)
pts=[
    "LLM infers the entire synchronized multi-robot plan; code only verifies & simulates.",
    "Pure-by-default design keeps the “LLM alone” claim clean; assisted mode is the ablation.",
    "All 8 planned milestones closed in ~4 weeks; prototype grew into a full framework.",
    "Next: run experiments, analyze ablations, and write up the paper.",
]
text(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(2.6),
     [[("—  "+p, 16, RGBColor(0xD7,0xDE,0xE8), False)] for p in pts],
     space_after=12, line_spacing=1.1)
text(s, Inches(0.9), SH-Inches(0.9), Inches(11.5), Inches(0.5),
     [[("LLM-as-MR-BT-Planner  ·  progress review  ·  17 June 2026", 12, MUTED, False)]])

out = r"D:\llm_as_mr_bt_planner\LLM-as-MR-BT-Planner_Progress.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
