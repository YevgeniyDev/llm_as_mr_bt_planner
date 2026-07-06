"""Assemble the revised presentation addressing all 5 feedback points."""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

ROOT = r"D:\llm_as_mr_bt_planner"
FIG = os.path.join(ROOT, "scratch_figs")
SRC = os.path.join(ROOT, "LLM-MultiRobot-BT-Planner_1.pptx")
DST = os.path.join(ROOT, "LLM-MultiRobot-BT-Planner_2.pptx")

prs = Presentation(SRC)
SW = prs.slide_width
def IN(v): return Inches(v)

# ---------- helpers ----------
def remove_shape(sh):
    sh._element.getparent().remove(sh._element)

def img_size(path):
    iw, ih = Image.open(path).size
    return iw / ih

def add_img(slide, path, top, width, target_bottom=6.98):
    asp = img_size(path)
    h = width / asp
    if top + h > target_bottom:
        h = target_bottom - top
        width = h * asp
    left = (13.3333 - width) / 2
    slide.shapes.add_picture(path, IN(left), IN(top), IN(width), IN(h))
    return top + h

def set_text(shape, s):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    runs = p0.runs
    if runs:
        runs[0].text = s
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.text = s
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)

def find_header(slide):
    """Return (section_label, title, pagenum) shapes by geometry."""
    sec = title = pg = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = Emu(sh.top).inches; l = Emu(sh.left).inches
        if abs(t - 0.42) < 0.1 and l < 2:
            sec = sh
        elif abs(t - 0.72) < 0.15 and l < 2:
            title = sh
        elif t > 6.9 and l > 11:
            pg = sh
    return sec, title, pg

# ---------- 1. Replace cartoon diagrams on slides 5,6,7 ----------
slides = prs.slides

# Slide 5 (index4): Problem Formulation
s5 = slides[4]
for sh in list(s5.shapes):
    if sh.shape_type == 13:
        remove_shape(sh)
b = add_img(s5, os.path.join(FIG, "fig_problem.png"), top=1.5, width=11.7)
print("slide5 fig bottom", round(b,3))

# Slide 6 (index5): Method - two figures
s6 = slides[5]
for sh in list(s6.shapes):
    if sh.shape_type == 13:
        remove_shape(sh)
asp_t = img_size(os.path.join(FIG, "fig_method_top.png"))
asp_b = img_size(os.path.join(FIG, "fig_method_bottom.png"))
w6 = 11.7
ht = w6 / asp_t; hb = w6 / asp_b
left6 = (13.3333 - w6) / 2
s6.shapes.add_picture(os.path.join(FIG, "fig_method_top.png"), IN(left6), IN(1.5), IN(w6), IN(ht))
s6.shapes.add_picture(os.path.join(FIG, "fig_method_bottom.png"), IN(left6), IN(1.5 + ht + 0.15), IN(w6), IN(hb))
print("slide6 bottoms", round(1.5+ht,3), round(1.5+ht+0.15+hb,3))

# Slide 7 (index6): Evaluation Protocol
s7 = slides[6]
for sh in list(s7.shapes):
    if sh.shape_type == 13:
        remove_shape(sh)
b = add_img(s7, os.path.join(FIG, "fig_eval.png"), top=1.5, width=12.1)
print("slide7 fig bottom", round(b,3))

# ---------- 2. Expand the Results table (slide index7) ----------
sR = slides[7]
old_tbl = None
for sh in list(sR.shapes):
    if sh.shape_type == 19:
        old_tbl = sh
tl, tt = 0.55, Emu(old_tbl.top).inches  # keep original top
tw = 12.23
remove_shape(old_tbl)

headers = ["Method", "Valid-BT ↑", "Goal succ. ↑", "Assign. acc. ↑",
           "Sync err. ↓", "Deadlock ↓", "Corr. rounds ↓", "BT size", "Time (s) ↓"]
methods = ["Manual (expert)", "MRBTP", "LLM-as-BT-Planner", "LLM-MARS", "Proposed (ours)"]
ncol, nrow = len(headers), len(methods) + 1
gf = sR.shapes.add_table(nrow, ncol, IN(tl), IN(tt), IN(tw), IN(2.35))
tbl = gf.table
tbl.first_row = True
tbl.horz_banding = False

# column widths
tbl.columns[0].width = IN(2.45)
rest = (tw - 2.45) / (ncol - 1)
for c in range(1, ncol):
    tbl.columns[c].width = IN(rest)

NAVY = RGBColor(0x1F, 0x3A, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INKC = RGBColor(0x22, 0x2A, 0x36)
GREENL = RGBColor(0xE3, 0xF3, 0xEA)
GREEND = RGBColor(0x16, 0x6B, 0x49)
ROWA = RGBColor(0xFF, 0xFF, 0xFF)
ROWB = RGBColor(0xF3, 0xF6, 0xF9)

def style_cell(cell, text, bold=False, color=INKC, size=9, align=PP_ALIGN.CENTER, fill=None):
    cell.margin_left = IN(0.05); cell.margin_right = IN(0.05)
    cell.margin_top = IN(0.03); cell.margin_bottom = IN(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    tf = cell.text_frame; tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = "Calibri"; f.color.rgb = color

# header row
for c, h in enumerate(headers):
    style_cell(tbl.cell(0, c), h, bold=True, color=WHITE, size=9,
               align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER, fill=NAVY)
# body rows
for r, m in enumerate(methods, start=1):
    is_ours = (m == "Proposed (ours)")
    base_fill = GREENL if is_ours else (ROWB if r % 2 == 0 else ROWA)
    for c in range(ncol):
        val = m if c == 0 else "—"
        style_cell(tbl.cell(r, c), val,
                   bold=is_ours, color=(GREEND if is_ours else INKC),
                   size=9.5 if c == 0 else 9,
                   align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER,
                   fill=base_fill)

# caption under table
cap = sR.shapes.add_textbox(IN(tl), IN(tt + 2.45), IN(tw), IN(0.3))
ctf = cap.text_frame; ctf.word_wrap = True
cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.LEFT
cr = cp.add_run()
cr.text = "↑ higher is better  ·  ↓ lower is better  ·  all methods scored by the same validator + simulator  ·  values pending full LLM evaluation (in progress)"
cf = cr.font; cf.size = Pt(9); cf.italic = True; cf.name = "Calibri"; cf.color.rgb = RGBColor(0x5B, 0x64, 0x72)

print("table rebuilt:", nrow, "x", ncol)

# ---------- 3. Create 3 new slides ----------
BLANK = prs.slide_layouts[0]
# reference header shapes from slide 6 (Method)
ref_sec, ref_title, ref_pg = find_header(slides[5])
ref_bg = slides[1].element.find(qn('p:cSld')).find(qn('p:bg'))

def new_slide(section, title, fig_path):
    sl = prs.slides.add_slide(BLANK)
    # explicit white background
    if ref_bg is not None:
        cSld = sl.element.find(qn('p:cSld'))
        cSld.insert(0, copy.deepcopy(ref_bg))
    # copy header shapes
    for ref, txt in [(ref_sec, section), (ref_title, title), (ref_pg, "0")]:
        el = copy.deepcopy(ref._element)
        sl.shapes._spTree.append(el)
    # set texts (last 3 shapes appended)
    shs = list(sl.shapes)
    set_text(shs[-3], section)
    set_text(shs[-2], title)
    set_text(shs[-1], "0")  # renumbered later
    # figure
    add_img(sl, fig_path, top=1.55, width=12.05)
    return sl

new_slide("4. OUR METHOD", "LLM Skills & Prompt Engineering", os.path.join(FIG, "fig_skills.png"))
new_slide("4. OUR METHOD", "Failure Detection & Recovery", os.path.join(FIG, "fig_failure.png"))
new_slide("5. EXPERIMENTS", "MuJoCo Physics Testing", os.path.join(FIG, "fig_mujoco.png"))
print("added 3 new slides; total now", len(prs.slides._sldIdLst))

# ---------- 4. Reorder slides ----------
# current order idx: 0 title,1 about,2 intro,3 related,4 problem,5 method,
#                    6 experiments,7 results,8 conclusion,9 skills,10 failure,11 mujoco
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
order = [0, 1, 2, 3, 4, 5, 9, 10, 6, 11, 7, 8]
for i in order:
    sldIdLst.append(ids[i])
print("reordered")

# ---------- 5. Renumber page numbers ----------
for i, sl in enumerate(prs.slides):
    if i < 2:
        continue
    _, _, pg = find_header(sl)
    if pg is not None:
        set_text(pg, str(i))
print("renumbered")

prs.save(DST)
print("SAVED ->", DST)
EOF = None
